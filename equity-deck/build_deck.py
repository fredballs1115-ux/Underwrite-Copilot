#!/usr/bin/env python3
"""
Build the 2300 N Street, NW – Investment Overview (Confidential Equity Memorandum).

Re-skins the Lincoln Property Company "Investment Overview" template (supplied as a
PDF export populated with a different asset) and re-populates it entirely with content
transposed from the Eastdil Secured "ES Opinion of Value" (BOV, Jun-2026) for the
2300 N Street, NW recapitalization – including a faithful recreation of Eastdil's
comparable-sale and financing data-point pages.

Sources: Eastdil Secured BOV (Jun-2026); CoStar reports (7/1/2026); ownership email
(loan balance $100M).
"""

import os, tempfile
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
_TMP = tempfile.mkdtemp(prefix="deckimg_")

def fit_image(path, w_in, h_in):
    """Center-crop an image to the target aspect ratio; return a temp path."""
    img = Image.open(path).convert("RGB")
    iw, ih = img.size
    target = w_in / h_in
    if iw / ih > target:
        nw = int(ih * target); x0 = (iw - nw) // 2
        img = img.crop((x0, 0, x0 + nw, ih))
    else:
        nh = int(iw / target); y0 = (ih - nh) // 2
        img = img.crop((0, y0, iw, y0 + nh))
    out = os.path.join(_TMP, "c_%x.jpg" % (abs(hash((path, round(w_in, 3), round(h_in, 3)))),))
    img.save(out, quality=88)
    return out

def photo(s, name, x_in, y_in, w_in, h_in):
    """Place an asset image, center-cropped to fill the target rectangle."""
    s.shapes.add_picture(fit_image(os.path.join(ASSETS, name), w_in, h_in),
                         Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))

# ----------------------------------------------------------------------------- brand
GOLD      = RGBColor(0xF5, 0xB4, 0x00)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x26, 0x26, 0x26)
GREY_TXT  = RGBColor(0x59, 0x59, 0x59)
ROW_GREY  = RGBColor(0xED, 0xED, 0xED)
HDR_DARK  = RGBColor(0x22, 0x22, 0x22)   # dark header row on comp tables
HDR_BROWN = RGBColor(0x3D, 0x30, 0x00)   # dark text on gold
LINE_GREY = RGBColor(0xBF, 0xBF, 0xBF)
PURPLE    = RGBColor(0xB0, 0xA4, 0xC6)
PURPLE_LT = RGBColor(0xBE, 0xB2, 0xD2)
PANEL     = RGBColor(0x1A, 0x1A, 0x1A)
PH_GREY   = RGBColor(0xD9, 0xD9, 0xD9)
# Lincoln-themed comp-matrix palette (Eastdil layout, Lincoln colors)
GOLD_LT   = RGBColor(0xFB, 0xEF, 0xCB)   # pale-gold row-label column
GREY_MED  = RGBColor(0x8C, 0x8C, 0x8C)   # pipeline (non-closed) status band
ALT       = RGBColor(0xF2, 0xF2, 0xF2)
PHOTO_BG  = RGBColor(0xDD, 0xDD, 0xDD)

TITLE_FONT = "Calibri Light"
BODY_FONT  = "Calibri"
SERIF_FONT = "Georgia"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]
SRC = "Source: Eastdil Secured Opinion of Value (Jun-2026)."

# ----------------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def bg(s, color):
    rect(s, 0, 0, Inches(10), Inches(5.625), color)

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=None, line_spacing=None, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None: p.space_after = Pt(space_after)
        if line_spacing is not None: p.line_spacing = line_spacing
        for (t, sz, col, bold, ital, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital; r.font.name = fnt
    return tb

def Rn(t, sz, col=DARK, bold=False, ital=False, fnt=BODY_FONT):
    return (t, sz, col, bold, ital, fnt)

def footer(s, page_num):
    ln = s.shapes.add_connector(2, Inches(0.42), Inches(5.30), Inches(9.58), Inches(5.30))
    ln.line.color.rgb = LINE_GREY; ln.line.width = Pt(0.75)
    txt(s, Inches(0.42), Inches(5.34), Inches(4), Inches(0.25),
        [[Rn("Lincoln Property Company", 8.5, GREY_TXT)]])
    txt(s, Inches(5.5), Inches(5.34), Inches(4.08), Inches(0.25),
        [[Rn(f"Confidential: Not For Distribution      |      {page_num}", 8.5, GREY_TXT)]],
        align=PP_ALIGN.RIGHT)

def content_header(s, title, page_num):
    rect(s, Inches(0.42), Inches(0.14), Inches(1.45), Inches(0.22), GOLD)
    txt(s, Inches(0.40), Inches(0.30), Inches(9), Inches(0.75),
        [[Rn(title, 32, DARK, False, False, TITLE_FONT)]])
    footer(s, page_num)

def gold_band(s, y, left_top, left_sub):
    """Gold report band with left label + right Lincoln wordmark (matches template)."""
    rect(s, Inches(0.42), y, Inches(9.16), Inches(0.62), GOLD)
    txt(s, Inches(0.60), Emu(int(y)+Inches(0.07)), Inches(6.5), Inches(0.3),
        [[Rn(left_top, 12.5, HDR_BROWN, True)]])
    txt(s, Inches(0.60), Emu(int(y)+Inches(0.33)), Inches(6.5), Inches(0.25),
        [[Rn(left_sub, 10, HDR_BROWN, False, True)]])
    txt(s, Inches(7.4), Emu(int(y)+Inches(0.10)), Inches(2.0), Inches(0.4),
        [[Rn("Lincoln", 24, HDR_BROWN, True, False, SERIF_FONT)]], align=PP_ALIGN.RIGHT)

def set_cell(cell, text, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT,
             ital=False, font=BODY_FONT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for j, line in enumerate(text.split("\n")):
        pp = p if j == 0 else tf.add_paragraph()
        pp.alignment = align
        r = pp.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = ital; r.font.name = font

def no_style(table):
    for el in table._tbl.findall(qn('a:tblPr')):
        el.set('firstRow', '0'); el.set('bandRow', '0')

def comp_table(s, x, y, w, h, cols, widths, rows, fontsize=7.6):
    tshape = s.shapes.add_table(len(rows)+1, len(cols), x, y, w, h)
    t = tshape.table; no_style(t)
    for i, wd in enumerate(widths): t.columns[i].width = Inches(wd)
    for c, hh in enumerate(cols):
        set_cell(t.cell(0, c), hh, fontsize, WHITE, bold=True, fill=HDR_DARK,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        fill = ROW_GREY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            set_cell(t.cell(r, c), val, fontsize, DARK, fill=fill,
                     align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    return t

def eastdil_matrix(s, y_top, y_bot, columns, row_labels, label_w, data_fs,
                   name_fs, name_h, photo_h, status_h, photos=None):
    """Eastdil-style comp matrix: properties as columns; name/photo/status bands
    on top; light-blue row labels down the left. `columns` = list of dicts with
    keys name, status, closed(bool), vals(list aligned to row_labels)."""
    ncol = 1 + len(columns)
    nrow = 3 + len(row_labels)
    x = Inches(0.42); total_w = Inches(9.16)
    tshape = s.shapes.add_table(nrow, ncol, x, Inches(y_top), total_w, Inches(y_bot - y_top))
    t = tshape.table; no_style(t)
    t.columns[0].width = Inches(label_w)
    data_w = (9.16 - label_w) / len(columns)
    for c in range(1, ncol): t.columns[c].width = Inches(data_w)
    attr_h = (y_bot - y_top - name_h - photo_h - status_h) / len(row_labels)
    t.rows[0].height = Inches(name_h)
    t.rows[1].height = Inches(photo_h)
    t.rows[2].height = Inches(status_h)
    for r in range(3, nrow): t.rows[r].height = Inches(attr_h)
    # corner cells (blank) above the labels
    for r in range(3):
        set_cell(t.cell(r, 0), "", data_fs, WHITE, fill=WHITE)
    photos = photos or {}
    for ci, col in enumerate(columns, start=1):
        set_cell(t.cell(0, ci), col["name"], name_fs, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.CENTER)
        if ci in photos:
            set_cell(t.cell(1, ci), "", 6, GREY_TXT, fill=WHITE)
        else:
            set_cell(t.cell(1, ci), "[ photo ]", 6, GREY_TXT, fill=PHOTO_BG, align=PP_ALIGN.CENTER, ital=True)
        set_cell(t.cell(2, ci), col["status"], data_fs, HDR_BROWN if col["closed"] else WHITE, bold=True,
                 fill=GOLD if col["closed"] else GREY_MED, align=PP_ALIGN.CENTER)
    # overlay real photos onto their photo-row cells
    for ci, fname in photos.items():
        x_in = 0.42 + label_w + (ci - 1) * data_w
        photo(s, fname, x_in, y_top + name_h, data_w, photo_h)
    for ri, lab in enumerate(row_labels):
        r = 3 + ri
        set_cell(t.cell(r, 0), lab, data_fs, DARK, bold=True, fill=GOLD_LT)
        for ci, col in enumerate(columns, start=1):
            fill = ALT if ri % 2 == 1 else WHITE
            set_cell(t.cell(r, ci), col["vals"][ri], data_fs, DARK, fill=fill, align=PP_ALIGN.CENTER)
    return t

def bullets(s, x, y, w, h, items, size=12, gap=6, lh=1.02, lead_color=GOLD):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, body = it
            paras.append([Rn("▪  ", size, lead_color, True), Rn(lead, size, DARK, True),
                          Rn(body, size, DARK)])
        else:
            paras.append([Rn("▪  ", size, lead_color, True), Rn(it, size, DARK)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=lh)

# ============================================================ 1 – COVER
s = slide(); bg(s, BLACK)
photo(s, "cover_2300_N.jpg", 4.55, 0, 5.45, 5.625)
txt(s, Inches(0.55), Inches(0.95), Inches(3.8), Inches(0.9),
    [[Rn("Lincoln", 40, GOLD, True, False, SERIF_FONT)]])
txt(s, Inches(0.6), Inches(1.95), Inches(3.8), Inches(0.5),
    [[Rn("Investment Overview", 20, WHITE, False, False, TITLE_FONT)]])
txt(s, Inches(0.6), Inches(2.65), Inches(3.9), Inches(1.2),
    [[Rn("2300 N Street, NW", 38, WHITE, True, False, TITLE_FONT)]])
txt(s, Inches(0.62), Inches(3.85), Inches(3.8), Inches(0.7),
    [[Rn("Washington, DC", 18, WHITE)], [Rn("West End Submarket", 18, WHITE)]], space_after=2)
txt(s, Inches(0.62), Inches(4.65), Inches(3.9), Inches(0.3),
    [[Rn("Lincoln Property Company  |  Oaktree Capital Management", 11, WHITE)]])
txt(s, Inches(0.62), Inches(5.15), Inches(3.8), Inches(0.3),
    [[Rn("Confidential Equity Memorandum", 11, WHITE, False, True)]])

# ============================================================ 2 – CONTENTS
s = slide(); bg(s, BLACK)
txt(s, Inches(0.42), Inches(1.85), Inches(2.6), Inches(0.8),
    [[Rn("Contents", 40, GOLD, False, False, TITLE_FONT)]])
items = [("Property Overview", "03"), ("Tenancy", "08"), ("Market", "12")]
y = 1.9
for name, pg in items:
    txt(s, Inches(3.6), Inches(y), Inches(5), Inches(0.5),
        [[Rn(name, 22, WHITE, False, False, TITLE_FONT)]])
    txt(s, Inches(8.15), Inches(y), Inches(1.1), Inches(0.5),
        [[Rn(pg, 22, WHITE, False, False, TITLE_FONT)]])
    y += 0.72

# ============================================================ 3 – PROPERTY OVERVIEW
s = slide(); bg(s, WHITE)
content_header(s, "Property Overview", 3)
txt(s, Inches(0.42), Inches(1.05), Inches(5.45), Inches(0.95),
    [[Rn("2300 N Street is a 282,943 SF, 8-story Class A office building in Washington, DC's "
        "supply-constrained West End submarket. Built in 1986 and renovated in 2018, the asset is "
        "87.8% leased with 9.4 years of WALT and best-in-class amenities.", 12, DARK)]],
    line_spacing=1.05)
rows = [
    ("Size (SF):", "282,943", "Submarket:", "West End"),
    ("Stories:", "8", "% Leased:", "87.8%"),
    ("Year Built / Reno:", "1986 / 2018", "In-Place WALT:", "9.4 Years"),
    ("Typical Floor:", "38,391 SF", "Parking:", "300 spaces (0.98/1,000 SF)"),
    ("Building Class:", "Class A", "Zoning:", "MU-10"),
    ("Land Area:", "1.13 AC", "Metro:", "Foggy Bottom-GWU (9-min walk)"),
    ("Amenities:", "Fitness, Roof Terrace,\nConferencing, Restaurant", "Ownership:", "Lincoln Property Co. / Oaktree"),
]
tshape = s.shapes.add_table(len(rows)+1, 4, Inches(0.42), Inches(2.05), Inches(5.35), Inches(2.9))
table = tshape.table; no_style(table)
table.columns[0].width = Inches(1.45); table.columns[1].width = Inches(1.30)
table.columns[2].width = Inches(1.15); table.columns[3].width = Inches(1.55)
hc = table.cell(0, 0); hc.merge(table.cell(0, 3))
set_cell(hc, "Property Description", 12, HDR_BROWN, bold=True, fill=GOLD, align=PP_ALIGN.CENTER)
for i, (a, b, c, d) in enumerate(rows, start=1):
    fill = ROW_GREY if i % 2 == 0 else WHITE
    set_cell(table.cell(i, 0), a, 9.5, DARK, bold=True, fill=fill)
    set_cell(table.cell(i, 1), b, 9.5, DARK, fill=fill)
    set_cell(table.cell(i, 2), c, 9.5, DARK, bold=True, fill=fill)
    set_cell(table.cell(i, 3), d, 9.5, DARK, fill=fill)
photo(s, "2300_N_exterior.jpg", 6.05, 1.05, 3.53, 1.9)
photo(s, "2300_N_roofterrace.jpg", 6.05, 3.05, 3.53, 1.9)

# ============================================================ 4 – LOCATION OVERVIEW
s = slide(); bg(s, WHITE)
content_header(s, "Location Overview", 4)
if not os.path.exists(os.path.join(ASSETS, "location_map.png")):
    import runpy
    runpy.run_path(os.path.join(os.path.dirname(os.path.abspath(__file__)), "make_map.py"))
photo(s, "location_map.png", 0.42, 1.05, 5.75, 3.95)
rect(s, Inches(6.35), Inches(1.05), Inches(3.23), Inches(0.4), GOLD)
txt(s, Inches(6.35), Inches(1.10), Inches(3.23), Inches(0.3),
    [[Rn("Location Highlights", 13, HDR_BROWN, True)]], align=PP_ALIGN.CENTER)
bullets(s, Inches(6.35), Inches(1.62), Inches(3.25), Inches(3.4), [
    ("West End", " — DC's premier office submarket, bordering Georgetown, Foggy Bottom and the CBD."),
    ("Metro", " — Foggy Bottom-GWU (9-min walk), Dupont Circle (13-min), Farragut North (15-min)."),
    ("Walkability", " — Walk Score 100, Transit 90, Bike 90 (“Exceptionally friendly”)."),
    ("Hotels", " — Fairmont, Park Hyatt, Westin Georgetown, Ritz-Carlton and Marriott within 2–3 blocks."),
    ("Regional access", " — Reagan National 12-min drive; Union Station / L'Enfant ~6–7 min."),
], size=10, gap=9, lh=1.0)
txt(s, Inches(0.42), Inches(5.05), Inches(9), Inches(0.25),
    [[Rn("Location highlights per CoStar (7/1/2026). Map is a schematic for orientation, not to scale.",
        8, GREY_TXT, False, True)]])

# ============================================================ 5 – VALUE CONSIDERATIONS (BOV p19)
s = slide(); bg(s, WHITE)
content_header(s, "Value Considerations", 5)
vc = [
    ("Hold Period:  ", "A 5-year hold is consistent with target investment horizons for active office "
        "capital and aligns with available financing."),
    ("Return Metrics:  ", "With material near-term leasing and capital, investors are likely to solve for an "
        "18–20% levered IRR and a 12–14% stabilized return on cost."),
    ("Capital Projections:  ", "Prior ownership invested ~$22.0M ($76 PSF) in base-building and amenity "
        "upgrades; $1.3M ($5 PSF) of outstanding TI/LCs is credited at closing; ES assumes a buyer underwrites "
        "$3.0M ($10 PSF) of near-term capital following sale."),
    ("Market Rents:  ", "Eastdil is underwriting typical office rents of $55.00–63.50 FSG, building on recent "
        "leasing at the asset."),
    ("Stabilized Occupancy:  ", "Occupancy steps to 77% after Aspen's downsize and Plan International's 2027 "
        "termination; ~40,000 SF of leasing in the first three years stabilizes occupancy at 90%."),
    ("Real Estate Taxes:  ", "ES assumes a 20% decrease in assessed value following sale, to $104.6M ($370 PSF), "
        "growing at inflation thereafter."),
    ("Market Financing:  ", "The most efficient structure is a floating-rate execution — 65% Initial LTV with 65% "
        "Future Funding for leasing/capex at SOFR + 3.75%."),
    ("Exit Cap Rate:  ", "ES underwrote a 9.5% exit cap, a $100.9M ($357 PSF) residual; absent clarity on Jackson "
        "& Campbell, investors underwrite a vacate (~50 bps of residual impact)."),
]
txt(s, Inches(0.42), Inches(1.10), Inches(9.16), Inches(4.0),
    [[Rn("▪  ", 11, GOLD, True), Rn(lead, 11, DARK, True), Rn(body, 11, DARK)] for lead, body in vc],
    space_after=7, line_spacing=1.0)
txt(s, Inches(0.42), Inches(5.02), Inches(9), Inches(0.25), [[Rn(SRC, 8, GREY_TXT, False, True)]])

# ============================================================ 6 – VALUATION (BOV p20)
s = slide(); bg(s, WHITE)
content_header(s, "Valuation", 6)
txt(s, Inches(0.42), Inches(1.05), Inches(9.15), Inches(0.5),
    [[Rn("Levered-IRR sensitivity — 5-year hold beginning October 2026; 65% LTV / 65% future funding "
        "at SOFR + 3.75% (7.43% coupon).", 12, DARK)]], line_spacing=1.03)
exit_caps = ["9.00%", "9.25%", "9.50%", "9.75%", "10.00%"]
price_rows = [
    ("$75,000,000", ["17.51%", "16.37%", "15.24%", "14.12%", "13.01%"]),
    ("$72,500,000", ["19.33%", "18.22%", "17.12%", "16.03%", "14.96%"]),
    ("$70,000,000", ["21.19%", "20.10%", "19.03%", "17.97%", "16.93%"]),
    ("$67,500,000", ["23.08%", "22.02%", "20.98%", "19.95%", "18.93%"]),
    ("$65,000,000", ["25.02%", "23.99%", "22.96%", "21.96%", "20.97%"]),
]
tshape = s.shapes.add_table(len(price_rows)+2, len(exit_caps)+1, Inches(0.42), Inches(1.95), Inches(6.0), Inches(2.4))
t = tshape.table; no_style(t)
t.columns[0].width = Inches(1.5)
for c in range(1, len(exit_caps)+1): t.columns[c].width = Inches(0.9)
top = t.cell(0, 1); top.merge(t.cell(0, len(exit_caps)))
set_cell(top, "Exit Cap Rate", 10.5, HDR_BROWN, bold=True, fill=GOLD, align=PP_ALIGN.CENTER)
set_cell(t.cell(0, 0), "Levered IRR", 9.5, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.CENTER)
set_cell(t.cell(1, 0), "Price", 9.5, DARK, bold=True, fill=ROW_GREY, align=PP_ALIGN.CENTER)
for c, cap in enumerate(exit_caps, start=1):
    set_cell(t.cell(1, c), cap, 9.5, DARK, bold=True, fill=ROW_GREY, align=PP_ALIGN.CENTER)
for r, (price, vals) in enumerate(price_rows, start=2):
    set_cell(t.cell(r, 0), price, 9.5, DARK, bold=True, fill=WHITE, align=PP_ALIGN.CENTER)
    for c, v in enumerate(vals, start=1):
        hl = (price == "$70,000,000" and exit_caps[c-1] == "9.50%")
        set_cell(t.cell(r, c), v, 9.5, HDR_BROWN if hl else DARK, bold=hl,
                 fill=GOLD if hl else WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(6.7), Inches(1.9), Inches(2.9), Inches(0.35),
    [[Rn("At ~$70.0M / $247 PSF", 13, DARK, True, False, TITLE_FONT)]])
metrics = [
    ("In-Place Cap Rate (87.8% occ)", "11.93%"), ("Adjusted Cap Rate (76.8% occ)", "8.86%"),
    ("Stabilized Return on Cost (Yr 5)", "12.04%"), ("Average Cash Yield", "9.42%"),
    ("Levered IRR", "19.03%"), ("Equity Multiple", "2.14x"), ("Reversion PSF (9.50% exit)", "$357"),
]
my = 2.35
for lab, val in metrics:
    txt(s, Inches(6.7), Inches(my), Inches(2.2), Inches(0.3), [[Rn(lab, 9.5, GREY_TXT)]])
    txt(s, Inches(8.9), Inches(my), Inches(0.7), Inches(0.3), [[Rn(val, 9.5, DARK, True)]], align=PP_ALIGN.RIGHT)
    my += 0.34
txt(s, Inches(0.42), Inches(4.55), Inches(6.0), Inches(0.7),
    [[Rn("Loan balance of $100M implies a ~25–35% discount to par at the ES value range. " + SRC,
        8, GREY_TXT, False, True)]], line_spacing=1.0)

# ============================================================ 7 – OPERATING CASH FLOW (BOV p21)
s = slide(); bg(s, WHITE)
content_header(s, "Operating Cash Flow", 7)
txt(s, Inches(0.42), Inches(1.05), Inches(9.15), Inches(0.5),
    [[Rn("Eastdil Secured five-year projection (years ending September 30). In-place NOI at "
        "10/1/2026 is $8.35M; near-term dips reflect lease-up downtime and abatement before "
        "stabilization.", 12, DARK)]], line_spacing=1.03)
cf_cols = ["($)", "FY 2027", "FY 2028", "FY 2029", "FY 2030", "FY 2031"]
cf_rows = [
    ("Effective Gross Income", "13,455,822", "13,177,239", "13,974,199", "15,774,954", "16,704,339"),
    ("Total Operating Expenses", "(6,942,899)", "(6,538,241)", "(6,703,809)", "(6,904,521)", "(7,083,080)"),
    ("Net Operating Income", "6,468,929", "6,472,555", "7,032,831", "8,616,206", "9,367,363"),
    ("Total Capital (TI/LC + CapEx)", "(1,748,577)", "(4,070,547)", "(3,217,200)", "(71,398)", "(215,515)"),
    ("Net Cash Flow", "4,720,352", "2,402,008", "3,815,631", "8,544,808", "9,151,849"),
]
tshape = s.shapes.add_table(len(cf_rows)+1, len(cf_cols), Inches(0.42), Inches(2.0), Inches(9.16), Inches(2.5))
t = tshape.table; no_style(t)
t.columns[0].width = Inches(2.56)
for c in range(1, 6): t.columns[c].width = Inches(1.32)
for c, h in enumerate(cf_cols):
    set_cell(t.cell(0, c), h, 10, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
for r, row in enumerate(cf_rows, start=1):
    emph = row[0] in ("Net Operating Income", "Net Cash Flow")
    fill = GOLD if emph else (ROW_GREY if r % 2 == 0 else WHITE)
    for c, val in enumerate(row):
        set_cell(t.cell(r, c), val, 10, HDR_BROWN if emph else DARK, bold=emph, fill=fill,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
txt(s, Inches(0.42), Inches(4.75), Inches(9), Inches(0.4),
    [[Rn("Capital reflects ~$3.0M ($10 PSF) of building CapEx (front-loaded FY27–29) plus tenant "
        "TI/LCs. " + SRC, 8, GREY_TXT, False, True)]], line_spacing=1.0)

# ============================================================ 8 – TENANCY DIVIDER
def divider(title):
    s = slide(); bg(s, PURPLE)
    txt(s, Inches(1.1), Inches(0.6), Inches(9), Inches(4.5),
        [[Rn("Lincoln", 200, PURPLE_LT, True, False, SERIF_FONT)]], wrap=False, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.55), Inches(2.2), Inches(5), Inches(1.0),
        [[Rn(title, 44, BLACK, False, False, TITLE_FONT)]])
    return s
divider("Tenancy")

# ============================================================ 9 – TENANT OVERVIEW (BOV p18)
s = slide(); bg(s, WHITE)
content_header(s, "Tenant Overview", 9)
cols = ["Tenant", "Floors", "SF", "% NRA", "Expiration", "Gross Rent ($/SF)"]
data = [
    ("Aspen Institute ¹", "6-8", "91,395", "32.3%", "May-41", "$71.56"),
    ("BlackIvy / CINQCARE", "2 & 6", "25,720", "9.1%", "Aug-34", "$55.87"),
    ("Bezos", "4-5", "23,874", "8.4%", "May-32", "$28.31"),
    ("Jackson & Campbell ²", "3", "23,735", "8.4%", "Aug-35", "$51.47"),
    ("Two Birds", "1", "17,275", "6.1%", "Jun-38", "$52.53"),
    ("Holtzman Vogel", "6", "16,722", "5.9%", "Nov-36", "$57.32"),
    ("Plan International ³", "3", "15,485", "5.5%", "Oct-27", "$67.03"),
    ("Compass Point", "4", "10,034", "3.5%", "Sep-37", "$53.50"),
    ("Capital Power", "5", "6,562", "2.3%", "Jun-31", "$55.00"),
    ("Bar Angie", "1", "4,435", "1.6%", "Mar-40", "$44.25"),
    ("Total Leased", "–", "248,300", "87.8%", "9.4 yrs", "$58.26"),
]
tshape = s.shapes.add_table(len(data)+2, len(cols), Inches(0.42), Inches(1.05), Inches(5.55), Inches(3.9))
t = tshape.table; no_style(t)
for i, wd in enumerate([1.55, 0.7, 0.9, 0.7, 0.95, 1.15]): t.columns[i].width = Inches(wd)
ban = t.cell(0, 0); ban.merge(t.cell(0, len(cols)-1))
set_cell(ban, "Tenant Summary", 11, HDR_BROWN, bold=True, fill=GOLD, align=PP_ALIGN.CENTER)
for c, h in enumerate(cols):
    set_cell(t.cell(1, c), h, 9, DARK, bold=True, fill=WHITE, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
for r, row in enumerate(data, start=2):
    total = row[0] == "Total Leased"
    fill = GOLD if total else (ROW_GREY if r % 2 == 0 else WHITE)
    for c, val in enumerate(row):
        set_cell(t.cell(r, c), val, 8.7, HDR_BROWN if total else DARK, bold=total, fill=fill,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
txt(s, Inches(6.2), Inches(1.05), Inches(3.4), Inches(0.35),
    [[Rn("Key Considerations", 13, DARK, True, False, TITLE_FONT)]])
bullets(s, Inches(6.2), Inches(1.45), Inches(3.4), Inches(3.4), [
    "88% leased today, stepping to 77% after Aspen's 6th-floor downsize — nearly 10 years of WALT.",
    "65,000+ SF (23% of NRA) of new leasing since Q4 2024; ~$1.3M of TI/LCs credited to a buyer at close.",
    "Anchor Aspen Institute extended through May-2041 across floors 6–8.",
    "Roll risk to watch: Plan International termination (Nov-2027) and Jackson & Campbell sublease "
    "marketing ahead of its Aug-2035 expiry.",
], size=10.5, gap=8, lh=1.02)
txt(s, Inches(0.42), Inches(5.02), Inches(9), Inches(0.25),
    [[Rn("¹ Gives back 15,536 SF Dec-2026, extends through May-2041.   ² 100% of space on sublease market; "
        "ES assumes vacate at expiry.   ³ Subleased; termination effective Nov-2027.", 7.5, GREY_TXT, False, True)]])

# ============================================================ 10 – RECENT LEASING ACTIVITY (BOV p23)
s = slide(); bg(s, WHITE)
content_header(s, "Recent Leasing Activity", 10)
gold_band(s, Inches(1.02), "2300 N STREET, NW", "Executed Leasing — Trailing 24 Months")
lcols = ["Tenant", "Suite", "NRA", "% NRA", "Commence", "Rent ($/SF)", "Term", "Free Rent", "TI ($/SF)", "Reimb."]
lrows = [
    ("Aspen Institute (Renewal)", "0700-0800", "75,859", "26.8%", "1/1/2027", "$63.50", "14 Yr 5 Mo", "9 Mo", "$0.00", "FSG"),
    ("Compass Point (New)", "0450-0460", "10,034", "3.5%", "6/1/2026", "$53.50", "11 Yr 4 Mo", "16 Mo", "$100.00", "FSG"),
    ("Capital Power (New)", "0520", "6,562", "2.3%", "1/1/2026", "$55.00", "5 Yr 6 Mo", "6 Mo", "$0.00", "FSG"),
    ("Bar Angie Storage (Exp.)", "0210", "2,841", "1.0%", "6/20/2025", "$25.34", "8 Yr 10 Mo", "0 Mo", "$0.00", "NNN"),
    ("Jackson & Campbell (Renewal)", "0300", "23,735", "8.4%", "3/1/2025", "$50.00", "10 Yr 6 Mo", "0 Mo", "$0.00", "FSG"),
    ("Holtzman Vogel (Exp.)", "0610", "10,280", "3.6%", "3/1/2025", "$55.50", "12 Yr", "20 Mo", "$135.00", "FSG"),
    ("Bezos (New)", "0400-0500", "23,874", "8.4%", "2/1/2025", "$27.35", "7 Yr 4 Mo", "3 Mo", "$0.00", "FSG"),
    ("Global Impact (New)", "0501A", "4,039", "1.4%", "2/1/2025", "$44.00", "7 Yr 4 Mo", "2 Mo", "$0.00", "FSG"),
    ("Timmons (New)", "0410B", "3,755", "1.3%", "2/1/2025", "$46.51", "1 Yr 9 Mo", "0 Mo", "$0.00", "FSG"),
    ("Bar Angie (New)", "0100", "4,435", "1.6%", "10/1/2024", "$41.71", "15 Yr 6 Mo", "6 Mo", "$281.85", "NNN"),
    ("Total / Avg.", "", "165,414", "58.5%", "", "$52.80", "139 Mo", "7.2 Mo", "$22.01", ""),
]
tshape = s.shapes.add_table(len(lrows)+1, len(lcols), Inches(0.42), Inches(1.75), Inches(9.16), Inches(3.35))
t = tshape.table; no_style(t)
lw = [2.05, 0.85, 0.75, 0.6, 0.85, 0.85, 0.95, 0.72, 0.82, 0.72]
for i, wd in enumerate(lw): t.columns[i].width = Inches(wd)
for c, h in enumerate(lcols):
    set_cell(t.cell(0, c), h, 8, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
for r, row in enumerate(lrows, start=1):
    total = row[0].startswith("Total")
    fill = GOLD if total else (ROW_GREY if r % 2 == 0 else WHITE)
    for c, val in enumerate(row):
        set_cell(t.cell(r, c), val, 7.8, HDR_BROWN if total else DARK, bold=total, fill=fill,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
txt(s, Inches(0.42), Inches(5.05), Inches(9), Inches(0.25), [[Rn(SRC + "  Rents shown gross.", 8, GREY_TXT, False, True)]])

# ============================================================ 11 – LEASING ASSUMPTIONS & CAPITAL (BOV p22-23)
s = slide(); bg(s, WHITE)
content_header(s, "Market Leasing Assumptions & Capital", 11)
# left: market leasing assumptions
txt(s, Inches(0.42), Inches(1.05), Inches(5.0), Inches(0.3),
    [[Rn("Market Leasing Assumptions", 13, DARK, True, False, TITLE_FONT)]])
la_cols = ["Space Type", "Rent", "Renew", "Down", "Term", "TI (N/R)"]
la_rows = [
    ("Office (Floors 7-8)", "$63.50", "70%", "18 mo", "11 Yr", "$120/$80"),
    ("Office (Floor 6)", "$59.00", "70%", "18 mo", "11 Yr", "$120/$80"),
    ("Office (Floors 4-5)", "$57.50", "70%", "18 mo", "11 Yr", "$120/$80"),
    ("Office (Floor 3)", "$55.00", "70%", "18 mo", "11 Yr", "$120/$80"),
    ("Town Hall Suites (4-5)", "$54.00", "70%", "12 mo", "8 Yr", "$75/$25"),
    ("Spec Suite (Floor 6)", "$52.00", "70%", "12 mo", "5.7 Yr", "$50/$35"),
    ("$36 Gross (Fl 2, 4-6)", "$36.00", "70%", "18 mo", "8 Yr", "$100/$30"),
    ("Retail", "$45.00", "75%", "12 mo", "10 Yr", "$95/$20"),
]
tshape = s.shapes.add_table(len(la_rows)+1, len(la_cols), Inches(0.42), Inches(1.45), Inches(5.2), Inches(3.1))
t = tshape.table; no_style(t)
for i, wd in enumerate([1.75, 0.72, 0.62, 0.62, 0.65, 0.84]): t.columns[i].width = Inches(wd)
for c, h in enumerate(la_cols):
    set_cell(t.cell(0, c), h, 8.5, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
for r, row in enumerate(la_rows, start=1):
    fill = ROW_GREY if r % 2 == 0 else WHITE
    for c, val in enumerate(row):
        set_cell(t.cell(r, c), val, 8.3, DARK, fill=fill, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
txt(s, Inches(0.42), Inches(4.55), Inches(5.2), Inches(0.5),
    [[Rn("Global: 2.5% inflation, 3.0% mgmt fee, $0.25 PSF reserves, 2.5% in-term bumps. " + SRC,
        7.5, GREY_TXT, False, True)]], line_spacing=1.0)
# right: capital plan
txt(s, Inches(5.9), Inches(1.05), Inches(3.7), Inches(0.3),
    [[Rn("Capital Plan", 13, DARK, True, False, TITLE_FONT)]])
cap_rows = [
    ("Prior Ownership Renovation", "$21.5M", "$76"),
    ("   Roof Terrace", "$4.6M", "$16"),
    ("   Common Areas", "$2.3M", "$8"),
    ("   Main Lobby / AC Units", "$4.0M", "$14"),
    ("   Fitness / Elevators / Garage", "$4.0M", "$14"),
    ("   Other base-building", "$6.6M", "$23"),
    ("Future Capital Plan (2026-28)", "$3.0M", "$11"),
]
tshape = s.shapes.add_table(len(cap_rows)+1, 3, Inches(5.9), Inches(1.45), Inches(3.68), Inches(2.7))
t = tshape.table; no_style(t)
for i, wd in enumerate([2.48, 0.7, 0.5]): t.columns[i].width = Inches(wd)
for c, h in enumerate(["Project", "Amount", "PSF"]):
    set_cell(t.cell(0, c), h, 8.5, WHITE, bold=True, fill=HDR_DARK, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
for r, row in enumerate(cap_rows, start=1):
    tot = row[0].startswith("Prior") or row[0].startswith("Future")
    fill = GOLD if tot else (ROW_GREY if r % 2 == 0 else WHITE)
    for c, val in enumerate(row):
        set_cell(t.cell(r, c), val, 8.3, HDR_BROWN if tot else DARK, bold=tot, fill=fill,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
txt(s, Inches(5.9), Inches(4.25), Inches(3.7), Inches(0.7),
    [[Rn("RE taxes: ES assumes a 20% assessment reduction post-sale to $104.6M ($370 PSF), "
        "growing at inflation. $1.3M of outstanding TI/LCs credited at close.", 7.5, GREY_TXT, False, True)]],
    line_spacing=1.0)

# ============================================================ 12 – MARKET DIVIDER
divider("Market")

# ============================================================ 13 – DC METRO OFFICE MARKET (BOV p5)
s = slide(); bg(s, WHITE)
content_header(s, "Washington, DC Metro Office Market", 13)
rect(s, Inches(0.42), Inches(1.05), Inches(4.5), Inches(0.4), GOLD)
txt(s, Inches(0.42), Inches(1.10), Inches(4.5), Inches(0.3), [[Rn("The Good", 15, HDR_BROWN, True)]], align=PP_ALIGN.CENTER)
rect(s, Inches(5.08), Inches(1.05), Inches(4.5), Inches(0.4), HDR_DARK)
txt(s, Inches(5.08), Inches(1.10), Inches(4.5), Inches(0.3), [[Rn("What's Challenging", 15, WHITE, True)]], align=PP_ALIGN.CENTER)
good = [
    "Office trade volume is accelerating — 3rd-highest 2025 US volume and 2nd-highest number of transactions.",
    "The return-to-office mandate has increased DC activity and added contractor demand.",
    "Pricing bands for office have formed.",
    "Newer-vintage product outperforms — nearly all net absorption is in post-2010 built/renovated space.",
    "Record rents are being achieved for new construction and high-quality space.",
    "Minimal new supply plus 6.9 MSF of office→residential/hotel conversions will keep improving fundamentals.",
]
chall = [
    "DC headlines (DOGE, GSA footprint, tariffs, National Guard, shutdown) create hurdles for investors.",
    "Headline vacancy above 20% in most submarkets is forcing heightened concessions, especially for commodity product.",
    "Leasing velocity remains mixed outside of Class-A space.",
    "Investors are underwriting significant capital amid escalating costs for non-trophy assets.",
]
bullets(s, Inches(0.42), Inches(1.6), Inches(4.45), Inches(3.4), good, size=10.5, gap=7, lh=1.0)
bullets(s, Inches(5.08), Inches(1.6), Inches(4.5), Inches(3.4), chall, size=10.5, gap=9, lh=1.0, lead_color=HDR_DARK)
txt(s, Inches(0.42), Inches(5.05), Inches(9), Inches(0.25), [[Rn(SRC, 8, GREY_TXT, False, True)]])

# ============================================================ 14 – RECENT DATA POINTS: CLASS-A (BOV p14)
s = slide(); bg(s, WHITE)
content_header(s, "Recent Data Points", 14)
gold_band(s, Inches(1.00), "WASHINGTON, DC — CLASS-A SALES", "Recent & Active Comparable Transactions")
def CA(name, status, closed, sz, lsd, walt, cap, price, buyer, seller, notes):
    return {"name": name, "status": status, "closed": closed,
            "vals": [sz, lsd, walt, cap, price, buyer, seller, notes]}
classA = [
    CA("888 16th", "Bidding", False, "107,000", "95%", "8.6 Yr", "8.50%", "$75.0M\n$700 PSF", "TBD", "Trammell Crow /\nMeadow", "Condo Regime"),
    CA("Hamilton Square", "Awarded", False, "248,000", "100%", "6.6 Yr", "8.90%", "$127.5M\n$514 PSF", "Confidential", "CommonWealth /\nCalPERS", "–"),
    CA("799 9th", "Awarded", False, "205,000", "82%", "6.6 Yr", "9.75%", "$100.0M\n$488 PSF", "Confidential", "Brookfield /\nAussieSuper", "–"),
    CA("1330 Connecticut", "Awarded", False, "259,000", "80%", "11.1 Yr", "9.00%", "$92.5M\n$358 PSF", "Confidential", "BXP", "–"),
    CA("2001 M", "Closed Jan-26", True, "284,000", "93%", "9.0 Yr", "8.90%", "$163.3M\n$577 PSF", "Stream Realty", "Brookfield", "Market Debt\n(Debt Fund)"),
    CA("1401 New York", "Closed Jan-26", True, "211,000", "94%→81%", "4.7 Yr", "11.50→9.00%", "$85.0M\n$402 PSF", "Carr / Barings", "AXA", "Market Debt\n(Debt Fund)"),
    CA("901 K", "Closed Jun-25", True, "220,000", "68%", "6.3 Yr", "8.00%", "$85.0M\n$386 PSF", "Shorenstein", "Carr Properties", "Market Debt\n(Debt Fund)"),
    CA("District Wharf\n(670/680/1000 Maine)", "Closed May-25", True, "750,000", "92%", "13.0 Yr", "6.50%", "$628.0M\n$835 PSF", "PSP", "Hoffman /\nMadison", "Market Debt\n(CMBS)"),
    CA("District Wharf\n(800 Maine)", "Closed May-25", True, "150,000", "100%", "8.5 Yr", "8.00%", "$63.2M\n$425 PSF", "PSP", "Hoffman /\nMadison", "Market Debt\n(CMBS)"),
    CA("2000 K", "Closed Jul-24", True, "233,000", "88%", "7.0 Yr", "8.00%", "$140.2M\n$600 PSF", "Spear Street", "Tishman /\nDeka", "Seller\nFinancing"),
    CA("1099 New York", "Closed Mar-24", True, "181,000", "95%", "9.0 Yr", "9.25%", "$95.0M\n$525 PSF", "FarmView /\nQuadrangle", "Credit Suisse", "Market Debt\n(CMBS)"),
]
eastdil_matrix(s, 1.72, 5.12, classA,
               ["Size (SF)", "% Leased", "WALT", "Cap Rate", "Price (PSF)", "Buyer", "Seller", "Notes"],
               label_w=0.98, data_fs=6.2, name_fs=6.5, name_h=0.48, photo_h=0.56, status_h=0.24,
               photos={1: "comp_888_16th.jpg", 2: "comp_Hamilton_Square.jpg",
                       3: "comp_799_9th.jpg", 4: "comp_1330_Connecticut.jpg",
                       5: "comp_2001_M.jpg", 6: "comp_1401_New_York.jpg",
                       7: "comp_901_K.jpg", 8: "comp_DW_670_Maine.jpg",
                       9: "comp_DW_800_Maine.jpg", 10: "comp_2000_K.jpg",
                       11: "comp_1099_New_York.jpg"})

# ============================================================ 15 – RECENT DATA POINTS: COMMODITY (BOV p15)
s = slide(); bg(s, WHITE)
content_header(s, "Recent Data Points", 15)
gold_band(s, Inches(1.00), "WASHINGTON, DC — COMMODITY SALES", "Recent & Active Comparable Transactions")
def CM(name, status, closed, sz, yr, lsd, walt, cap, price, buyer, seller):
    return {"name": name, "status": status, "closed": closed,
            "vals": [sz, yr, lsd, walt, cap, price, buyer, seller]}
commodity = [
    CM("700 6th", "Bidding", False, "306,000", "2009", "79%", "7.8 Yr", "11.8%", "$85M\n$278", "TBD", "Principal\n(Lender)"),
    CM("1350 Eye", "Bidding", False, "415,861", "1989/14", "75%", "8.0 Yr", "10.0%", "$130M\n$348", "TBD", "MetLife"),
    CM("1775 Penn", "Awarded", False, "156,482", "1975/09", "41%", "5.5 Yr", "12.8%\n(ROC)", "$37.5M\n$240", "Confid.", "Woodwind"),
    CM("Sumner\nSquare", "Under\nContract", False, "242,000", "1985", "100%", "TBD", "TBD", "$63.0M\n$261", "Vireon", "BXP"),
    CM("2445 M", "Closed\nApr-26", True, "300,000", "1986/19", "80→75%", "7.7 Yr", "9.0→8.5%", "$101.0M\n$337", "Eagle Cliff", "Mesa West"),
    CM("Watergate", "Closed\nMar-26", True, "309,000", "1972/02", "93→20%", "N/A", "N/A", "$52.5M\n$170", "JetSet", "Elme"),
    CM("1400 K", "Closed\nOct-25", True, "200,000", "1982/24", "65%", "5.3 Yr", "9.5%", "$35.5M\n$180", "Taicoon", "Brookfield"),
    CM("2033 K", "Closed\nJul-25", True, "128,767", "1975/19", "38%", "–", "N/A", "$20.3M\n$170", "In-Rel", "ARA"),
    CM("1325 G", "Closed\nJun-25", True, "305,259", "1969/04", "62%", "2.5 Yr", "TBD", "$25.0M\n$82", "Farmview", "US Bank"),
    CM("1750 H", "Closed\nJun-25", True, "124,000", "2002", "55%", "3.8 Yr", "TBD", "$28.8M\n$232", "Douglas", "State Farm"),
    CM("300 M", "Closed\nJun-25", True, "282,354", "2001", "50%", "3.8 Yr", "11.5%", "$28.0M\n$100", "Confid.", "Potomac"),
    CM("2101 L", "Closed\nDec-24", True, "375,000", "1975/07", "75%", "8.0 Yr", "10.0%", "$110.0M\n$293", "ELV", "JBGS /\nNW Mutual"),
    CM("1250 H", "Closed\nNov-24", True, "196,490", "1992", "53%", "7.4 Yr", "10.7%", "$27.5M\n$137", "Banyan", "Equity CW"),
    CM("2001 L", "Closed\nOct-24", True, "172,135", "1986/14", "66%", "6.3 Yr", "10.7%", "$30.55M\n$180", "Melrose", "AEW /\nSwiss RE"),
    CM("1899 L", "Closed\nMar-24", True, "150,502", "1978/22", "56%", "4.5 Yr", "9.5%", "$26.65M\n$180", "Taicoon", "BlackRock"),
]
eastdil_matrix(s, 1.72, 5.12, commodity,
               ["Size (SF)", "Year Built", "% Leased", "WALT", "Cap Rate", "Price (PSF)", "Buyer", "Seller"],
               label_w=0.92, data_fs=5.6, name_fs=6.0, name_h=0.46, photo_h=0.46, status_h=0.32,
               photos={1: "comp_700_6th.jpg", 3: "comp_1775_Penn.jpg",
                       4: "comp_Sumner_Square.jpg", 5: "comp_2445_M.jpg",
                       6: "comp_Watergate.jpg"})

# ============================================================ 16 – SELECT ES OFFICE FINANCINGS (BOV p16)
s = slide(); bg(s, WHITE)
content_header(s, "Recent Data Points", 16)
gold_band(s, Inches(1.00), "WASHINGTON, DC METRO — DEBT DATA POINTS", "Recent Eastdil Secured Office Financings")
def FN(name, sponsor, status, closed, sf, walt, lsd, loan, loanpsf, ltv, dy, lender, pricing):
    return {"name": name, "status": status, "closed": closed,
            "vals": [sponsor, sf, walt, lsd, loan, loanpsf, ltv, dy, lender, pricing]}
fins = [
    FN("Reston Metro\nPlaza", "Comstock", "Marketing", False, "769,830", "8.3 Yr", "99%", "$299.1M", "$389", "75% LTV", "9.6%", "TBD", "TBD"),
    FN("799 9th\nSt, NW", "Jemal Equities", "Awarding", False, "203,329", "6.5 Yr", "82%", "$83.1M", "$409", "75% LTC", "13.4%", "Debt Fund", "SOFR +\n~4.25%"),
    FN("Market\nSquare", "Elliott / MC /\nPRP", "Under\nContract", False, "694,436", "5.3 Yr", "90%", "$340.0M", "$490", "70% LTV", "12.0%", "SASB", "SOFR +\n~2.75%"),
    FN("2001 M", "Stream Realty", "Closed\nJan-26", True, "283,190", "9.0 Yr", "93%", "$118.5M", "$418", "69% LTC\n(68%I/75%FF)", "13.3%", "Debt Fund", "SOFR +\n2.90%"),
    FN("1401\nNew York", "Carr / Barings", "Closed\nJan-26", True, "201,976", "4.7 Yr", "94→81%", "$73.4M", "$344", "68% LTC\n(65%I/84%FF)", "13.9%", "Debt Fund", "SOFR +\n4.05%"),
    FN("Three-Asset\nPortfolio", "Carr Properties", "Closed\nJul-25", True, "791,098", "9.3 Yr", "89%", "$278.3M", "$352", "60% LTV", "12.0%", "Bank\nBalance Sheet", "SOFR +\n3.25%"),
    FN("901 K St", "Shorenstein", "Closed\nJun-25", True, "220,000", "6.3 Yr", "68%", "$76.1M", "$346", "67.5% LTC\n(60%I/100%FF)", "12.6%", "Debt Fund", "SOFR +\n4.15%"),
    FN("Boro Tower", "Rockefeller /\nTMG", "Closed\nMay-25", True, "436,795", "7.3 Yr", "98→88%", "$175.0M", "$401", "70% LTV", "11.5→10.0%", "Debt Fund", "SOFR + 4.75%\n(+4.40% b/d)"),
]
eastdil_matrix(s, 1.72, 5.12, fins,
               ["Sponsor", "SF", "WALT", "% Leased", "Loan Amount", "Loan PSF", "LTV / LTC", "In-Place DY", "Lender", "Pricing"],
               label_w=1.28, data_fs=6.0, name_fs=6.6, name_h=0.40, photo_h=0.40, status_h=0.28,
               photos={2: "comp_799_9th.jpg", 4: "comp_2001_M.jpg",
                       5: "comp_1401_New_York.jpg", 7: "comp_901_K.jpg"})

# ============================================================ 17 – CONTACT
s = slide(); bg(s, BLACK)
txt(s, Inches(0.42), Inches(1.55), Inches(3), Inches(0.8),
    [[Rn("Contact", 40, GOLD, False, False, TITLE_FONT)]])
contacts = [
    ("Matt Nicholson", "Executive Vice President", "mnicholson@lpc.com"),
    ("Jake Guttman", "Director", "jguttman@lpc.com"),
    ("Jackson Coviello", "Assistant Vice President", "jcoviello@lpc.com"),
    ("Jared Hicks", "Analyst", "jahicks@lpc.com"),
]
for i, (nm, title, email) in enumerate(contacts):
    x = Inches(0.5 + i*2.35)
    txt(s, x, Inches(3.0), Inches(2.3), Inches(0.4), [[Rn(nm, 17, WHITE, False, False, TITLE_FONT)]])
    txt(s, x, Inches(3.42), Inches(2.3), Inches(0.3), [[Rn(title, 11, WHITE)]])
    txt(s, x, Inches(3.70), Inches(2.3), Inches(0.3), [[Rn(email, 11, WHITE)]])
txt(s, Inches(0.62), Inches(4.55), Inches(6), Inches(1.0),
    [[Rn("Lincoln Property Company", 10, WHITE)], [Rn("101 Constitution Ave. NW, Suite 325 East", 10, WHITE)],
     [Rn("Washington, DC 20001", 10, WHITE)], [Rn("www.lpc.com", 10, WHITE)]], space_after=3)

out = "/home/user/Underwrite-Copilot/equity-deck/2300_N_Street_NW_Investment_Overview.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
